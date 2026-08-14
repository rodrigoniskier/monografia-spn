import json
import logging
import re
from dataclasses import dataclass, field
from datetime import date, datetime
from io import BytesIO
from pathlib import Path
from urllib.parse import unquote, urlsplit

import requests
from bs4 import BeautifulSoup
from django.conf import settings
from docx import Document
from pptx import Presentation
from pypdf import PdfReader

from ..models import Reference
from .people import parse_people
from .security import RemoteFetchError, UnsafeURL, normalize_url, safe_fetch

logger = logging.getLogger(__name__)
DOI_RE = re.compile(r"\b10\.\d{4,9}/[-._;()/:A-Z0-9]+", re.I)
YEAR_RE = re.compile(r"\b(?:18|19|20)\d{2}\b")


@dataclass
class ExtractionResult:
    fields: dict = field(default_factory=dict)
    warnings: list[str] = field(default_factory=list)


def compact(value) -> str:
    if value is None:
        return ""
    return re.sub(r"\s+", " ", str(value)).strip()


def first_value(value):
    if isinstance(value, list):
        return value[0] if value else ""
    return value or ""


def find_doi(text: str) -> str:
    match = DOI_RE.search(text or "")
    return match.group(0).rstrip(".,;)]}") if match else ""


def find_year(value) -> str:
    if isinstance(value, (date, datetime)):
        return str(value.year)
    match = YEAR_RE.search(compact(value))
    return match.group(0) if match else ""


def plausible_title(lines: list[str], fallback: str) -> str:
    ignored = ("http://", "https://", "doi:", "abstract", "resumo", "issn", "isbn")
    for line in lines[:40]:
        value = compact(line)
        if (
            12 <= len(value) <= 500
            and not value.casefold().startswith(ignored)
            and not YEAR_RE.fullmatch(value)
        ):
            return value
    return fallback


def crossref_metadata(doi: str) -> dict:
    doi = find_doi(doi) or compact(doi)
    if not doi or not settings.CROSSREF_ENRICH:
        return {}
    headers = {
        "User-Agent": f"CitaRN/1.0 (mailto:{settings.CROSSREF_MAILTO or 'contato-nao-informado'})"
    }
    try:
        response = requests.get(
            f"https://api.crossref.org/works/{requests.utils.quote(doi, safe='')}",
            headers=headers,
            timeout=(3, 6),
        )
        response.raise_for_status()
        message = response.json().get("message", {})
    except (requests.RequestException, ValueError) as exc:
        logger.info("Crossref enrichment unavailable for DOI %s: %s", doi, exc)
        return {}

    type_map = {
        "journal-article": Reference.Type.JOURNAL_ARTICLE,
        "book": Reference.Type.BOOK,
        "book-chapter": Reference.Type.BOOK_CHAPTER,
        "proceedings-article": Reference.Type.CONFERENCE,
        "dissertation": Reference.Type.THESIS,
        "report": Reference.Type.REPORT,
    }
    published = (
        message.get("published-print")
        or message.get("published-online")
        or message.get("issued")
        or {}
    )
    date_parts = published.get("date-parts", [[]])
    year = str(date_parts[0][0]) if date_parts and date_parts[0] else ""
    authors = [
        {
            key: compact(person.get(key))
            for key in ("family", "given", "literal")
            if compact(person.get(key))
        }
        for person in message.get("author", [])
    ]
    editors = [
        {
            key: compact(person.get(key))
            for key in ("family", "given", "literal")
            if compact(person.get(key))
        }
        for person in message.get("editor", [])
    ]
    return {
        "reference_type": type_map.get(
            message.get("type"), Reference.Type.JOURNAL_ARTICLE
        ),
        "title": compact(first_value(message.get("title"))),
        "subtitle": compact(first_value(message.get("subtitle"))),
        "authors": [person for person in authors if person],
        "editors": [person for person in editors if person],
        "year": year,
        "container_title": compact(first_value(message.get("container-title"))),
        "publisher": compact(message.get("publisher")),
        "publisher_place": compact(message.get("publisher-location")),
        "edition": compact(message.get("edition-number")),
        "volume": compact(message.get("volume")),
        "issue": compact(message.get("issue")),
        "pages": compact(message.get("page")),
        "doi": compact(message.get("DOI")) or doi,
        "url": compact(message.get("URL")),
        "language": compact(message.get("language")),
    }


def finalize(
    fields: dict, warnings: list[str], fallback_title: str
) -> ExtractionResult:
    fields = {
        key: value for key, value in fields.items() if value not in (None, "", [], {})
    }
    search_text = " ".join(
        compact(value) for value in fields.values() if isinstance(value, (str, int))
    )
    fields.setdefault("doi", find_doi(search_text))
    fields.setdefault("year", find_year(search_text))
    fields.setdefault("title", fallback_title)
    fields.setdefault("reference_type", Reference.Type.JOURNAL_ARTICLE)
    fields.setdefault("language", "pt")
    if fields.get("doi"):
        enriched = crossref_metadata(fields["doi"])
        if enriched:
            # A DOI identifica a obra de forma inequívoca; os dados do Crossref são preferidos.
            fields.update(
                {
                    key: value
                    for key, value in enriched.items()
                    if value not in (None, "", [], {})
                }
            )
    if not fields.get("authors"):
        warnings.append("Autores não identificados automaticamente.")
    if not fields.get("year"):
        warnings.append("Ano de publicação não identificado automaticamente.")
    warnings.append(
        "Revise os metadados antes de usar a referência em um trabalho acadêmico."
    )
    return ExtractionResult(fields=fields, warnings=list(dict.fromkeys(warnings)))


def extract_pdf_bytes(data: bytes, fallback_title: str) -> ExtractionResult:
    warnings = []
    try:
        reader = PdfReader(BytesIO(data))
        metadata = reader.metadata or {}
        page_text = []
        for page in reader.pages[:4]:
            try:
                page_text.append(page.extract_text() or "")
            except Exception:
                warnings.append("Parte do texto do PDF não pôde ser lida.")
        full_text = "\n".join(page_text)
        lines = [compact(line) for line in full_text.splitlines() if compact(line)]
        title = compact(getattr(metadata, "title", "")) or plausible_title(
            lines, fallback_title
        )
        author = compact(getattr(metadata, "author", ""))
        fields = {
            "title": title,
            "authors": parse_people(author),
            "year": find_year(getattr(metadata, "creation_date", ""))
            or find_year(full_text),
            "doi": find_doi(full_text),
        }
        return finalize(fields, warnings, fallback_title)
    except Exception as exc:
        logger.warning("PDF extraction failed: %s", exc)
        return ExtractionResult(
            fields={"title": fallback_title, "reference_type": Reference.Type.OTHER},
            warnings=[
                "O PDF foi salvo, mas seus metadados não puderam ser extraídos.",
                "Preencha os campos manualmente.",
            ],
        )


def extract_docx_bytes(data: bytes, fallback_title: str) -> ExtractionResult:
    try:
        document = Document(BytesIO(data))
        properties = document.core_properties
        lines = [
            compact(paragraph.text)
            for paragraph in document.paragraphs
            if compact(paragraph.text)
        ]
        text = "\n".join(lines[:150])
        fields = {
            "title": compact(properties.title)
            or plausible_title(lines, fallback_title),
            "authors": parse_people(compact(properties.author)),
            "year": find_year(properties.created) or find_year(text),
            "doi": find_doi(text),
            "language": compact(properties.language),
        }
        return finalize(fields, [], fallback_title)
    except Exception as exc:
        logger.warning("DOCX extraction failed: %s", exc)
        return ExtractionResult(
            fields={"title": fallback_title, "reference_type": Reference.Type.OTHER},
            warnings=[
                "O DOCX foi salvo, mas seus metadados não puderam ser extraídos.",
                "Preencha os campos manualmente.",
            ],
        )


def extract_pptx_bytes(data: bytes, fallback_title: str) -> ExtractionResult:
    try:
        presentation = Presentation(BytesIO(data))
        properties = presentation.core_properties
        lines = []
        for slide in presentation.slides[:6]:
            for shape in slide.shapes:
                if hasattr(shape, "text") and compact(shape.text):
                    lines.extend(compact(shape.text).splitlines())
        text = "\n".join(lines)
        fields = {
            "reference_type": Reference.Type.CONFERENCE,
            "title": compact(properties.title)
            or plausible_title(lines, fallback_title),
            "authors": parse_people(compact(properties.author)),
            "year": find_year(properties.created) or find_year(text),
            "doi": find_doi(text),
        }
        return finalize(fields, [], fallback_title)
    except Exception as exc:
        logger.warning("PPTX extraction failed: %s", exc)
        return ExtractionResult(
            fields={"title": fallback_title, "reference_type": Reference.Type.OTHER},
            warnings=[
                "O PPTX foi salvo, mas seus metadados não puderam ser extraídos.",
                "Preencha os campos manualmente.",
            ],
        )


def extract_file(uploaded_file) -> ExtractionResult:
    filename = Path(uploaded_file.name).name
    fallback_title = (
        Path(filename).stem.replace("_", " ").replace("-", " ").strip().title()
    )
    extension = Path(filename).suffix.lower()
    data = uploaded_file.read()
    uploaded_file.seek(0)
    if extension == ".pdf":
        return extract_pdf_bytes(data, fallback_title)
    if extension == ".docx":
        return extract_docx_bytes(data, fallback_title)
    if extension == ".pptx":
        return extract_pptx_bytes(data, fallback_title)
    raise ValueError("Formato de arquivo não suportado.")


def _meta_values(soup: BeautifulSoup) -> dict[str, list[str]]:
    values: dict[str, list[str]] = {}
    for tag in soup.find_all("meta"):
        key = compact(tag.get("name") or tag.get("property")).casefold()
        value = compact(tag.get("content"))
        if key and value:
            values.setdefault(key, []).append(value)
    return values


def _json_ld_records(soup: BeautifulSoup) -> list[dict]:
    records = []
    for script in soup.find_all("script", attrs={"type": "application/ld+json"}):
        try:
            payload = json.loads(script.string or script.get_text() or "{}")
        except (TypeError, json.JSONDecodeError):
            continue
        candidates = payload if isinstance(payload, list) else [payload]
        for candidate in candidates:
            if isinstance(candidate, dict) and isinstance(
                candidate.get("@graph"), list
            ):
                records.extend(
                    item for item in candidate["@graph"] if isinstance(item, dict)
                )
            elif isinstance(candidate, dict):
                records.append(candidate)
    return records


def extract_html(data: bytes, source_url: str) -> ExtractionResult:
    soup = BeautifulSoup(data, "html.parser")
    meta = _meta_values(soup)
    json_ld = _json_ld_records(soup)

    def meta_first(*keys):
        for key in keys:
            if meta.get(key):
                return meta[key][0]
        return ""

    title = meta_first(
        "citation_title", "dc.title", "dcterms.title", "og:title", "twitter:title"
    )
    if not title and soup.title:
        title = compact(soup.title.string)
    authors = (
        meta.get("citation_author", [])
        or meta.get("dc.creator", [])
        or meta.get("author", [])
    )
    fields = {
        "reference_type": Reference.Type.WEBSITE,
        "title": title,
        "authors": parse_people(authors),
        "year": find_year(
            meta_first(
                "citation_publication_date", "article:published_time", "date", "dc.date"
            )
        ),
        "container_title": meta_first("citation_journal_title", "og:site_name"),
        "publisher": meta_first("dc.publisher", "citation_publisher"),
        "volume": meta_first("citation_volume"),
        "issue": meta_first("citation_issue"),
        "pages": meta_first("citation_firstpage"),
        "doi": meta_first("citation_doi", "dc.identifier", "dc.identifier.doi"),
        "url": source_url,
        "source_url": source_url,
    }
    last_page = meta_first("citation_lastpage")
    if fields.get("pages") and last_page:
        fields["pages"] = f"{fields['pages']}-{last_page}"

    for record in json_ld:
        record_type = record.get("@type", "")
        types = record_type if isinstance(record_type, list) else [record_type]
        if not any(
            item in {"ScholarlyArticle", "Article", "NewsArticle", "Book", "Chapter"}
            for item in types
        ):
            continue
        fields["title"] = fields.get("title") or compact(
            record.get("headline") or record.get("name")
        )
        fields["year"] = fields.get("year") or find_year(record.get("datePublished"))
        raw_authors = record.get("author", [])
        if isinstance(raw_authors, dict):
            raw_authors = [raw_authors]
        structured = []
        for author in raw_authors if isinstance(raw_authors, list) else []:
            if isinstance(author, dict):
                if author.get("familyName"):
                    structured.append(
                        {
                            "family": compact(author.get("familyName")),
                            "given": compact(author.get("givenName")),
                        }
                    )
                elif author.get("name"):
                    structured.extend(parse_people(author.get("name")))
            elif author:
                structured.extend(parse_people(author))
        fields["authors"] = fields.get("authors") or structured
        fields["doi"] = fields.get("doi") or find_doi(compact(record.get("identifier")))
        if "Book" in types:
            fields["reference_type"] = Reference.Type.BOOK
        elif "Chapter" in types:
            fields["reference_type"] = Reference.Type.BOOK_CHAPTER
        else:
            fields["reference_type"] = Reference.Type.JOURNAL_ARTICLE
        break

    hostname = urlsplit(source_url).hostname or source_url
    return finalize(fields, [], hostname)


def extract_url(value: str) -> ExtractionResult:
    value = normalize_url(value)
    parsed = urlsplit(value)
    if parsed.hostname and parsed.hostname.casefold() in {"doi.org", "dx.doi.org"}:
        doi = unquote(parsed.path.lstrip("/"))
        enriched = crossref_metadata(doi)
        if enriched:
            enriched.setdefault("source_url", value)
            enriched.setdefault("url", value)
            return finalize(enriched, [], value)
    try:
        data, content_type, final_url = safe_fetch(value)
        if content_type == "application/pdf" or final_url.casefold().endswith(".pdf"):
            result = extract_pdf_bytes(
                data, Path(urlsplit(final_url).path).stem or "Documento em PDF"
            )
            result.fields["source_url"] = final_url
            result.fields.setdefault("url", final_url)
            return result
        return extract_html(data, final_url)
    except (UnsafeURL, RemoteFetchError) as exc:
        hostname = urlsplit(value).hostname or value
        return ExtractionResult(
            fields={
                "reference_type": Reference.Type.WEBSITE,
                "title": hostname,
                "url": value if value.startswith(("http://", "https://")) else "",
                "source_url": value
                if value.startswith(("http://", "https://"))
                else "",
            },
            warnings=[
                str(exc),
                "A URL foi cadastrada; complete os metadados manualmente.",
            ],
        )
